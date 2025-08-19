import os
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image

def px_to_emu(px, dpi=96):
    return int(px * 914400 / dpi)

def add_image_slide(prs, img_path, dpi=200):
    with Image.open(img_path) as im:
        img_w, img_h = im.size
    slide_width = px_to_emu(img_w, dpi)
    slide_height = px_to_emu(img_h, dpi)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        img_path,
        0, 0,
        width=slide_width,
        height=slide_height
    )

def convert_pdf_to_ppt(input_path, output_dir, timestr=""):
    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_filename = f"{base_name}_{timestr}.pptx" if timestr else f"{base_name}.pptx"
    output_path = os.path.join(output_dir, output_filename)

    dpi = 200
    images = convert_from_path(input_path, dpi=dpi, fmt="png")

    prs = Presentation()
    img_w, img_h = images[0].size
    prs.slide_width = px_to_emu(img_w, dpi)
    prs.slide_height = px_to_emu(img_h, dpi)

    tmp_dir = os.path.join(output_dir, "tmp_img")
    os.makedirs(tmp_dir, exist_ok=True)
    img_paths = []
    for idx, img in enumerate(images):
        img_file = os.path.join(tmp_dir, f"page_{idx+1}.png")
        img.save(img_file)
        img_paths.append(img_file)
        add_image_slide(prs, img_file, dpi)

    prs.save(output_path)

    for img_file in img_paths:
        os.remove(img_file)
    os.rmdir(tmp_dir)

    return output_path, output_filename